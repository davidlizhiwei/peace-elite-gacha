#!/usr/bin/env python3
"""
Create multiplication table in Word and PowerPoint formats
"""

import os
from docx import Document
from docx.shared import Pt, Inches, RGBColor as DocxRGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor as PptxRGBColor
from pptx.enum.text import PP_ALIGN

# Generate multiplication table data
def get_multiplication_table():
    """Generate 9x9 multiplication table"""
    table = []
    for i in range(1, 10):
        row = []
        for j in range(1, i + 1):
            row.append(f"{j}×{i}={i*j}")
        table.append(row)
    return table

# Create Word document
def create_word_doc():
    """Create Word document with multiplication table"""
    doc = Document()

    # Title
    title = doc.add_heading('乘法口诀表', 0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER

    # Add subtitle
    subtitle = doc.add_paragraph('九九乘法表')
    subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
    subtitle.runs[0].font.size = Pt(14)
    subtitle.runs[0].font.color.rgb = DocxRGBColor(100, 100, 100)

    # Create table
    table = doc.add_table(rows=9, cols=9)
    table.style = 'Table Grid'

    # Fill table with multiplication data
    mult_data = get_multiplication_table()

    for i in range(9):
        row = table.rows[i]
        for j in range(9):
            cell = row.cells[j]
            if j <= i:
                cell.text = mult_data[i][j]
                cell.paragraphs[0].alignment = WD_ALIGN_PARAGRAPH.CENTER
                run = cell.paragraphs[0].runs[0]
                run.font.size = Pt(12)
                run.font.bold = True
                # Color code by result
                result = (i + 1) * (j + 1)
                if result <= 10:
                    run.font.color.rgb = DocxRGBColor(0, 100, 0)  # Green for small numbers
                elif result <= 50:
                    run.font.color.rgb = DocxRGBColor(0, 0, 150)  # Blue for medium
                else:
                    run.font.color.rgb = DocxRGBColor(150, 0, 0)  # Red for large
            else:
                cell.text = ""

    # Add explanation
    doc.add_paragraph()
    explanation = doc.add_paragraph()
    explanation.add_run('说明：').bold = True
    explanation.add_run('这是中国传统的九九乘法表，共 9 行 9 列，从 1×1=1 到 9×9=81。')

    # Save document
    output_path = '/Users/davidli/lobsterai/project/乘法口诀表.docx'
    doc.save(output_path)
    print(f"✓ Word document saved to: {output_path}")
    return output_path

# Create PowerPoint presentation
def create_ppt():
    """Create PowerPoint with multiplication table"""
    prs = Presentation()

    # Title slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "乘法口诀表"
    subtitle.text = "九九乘法表\n中国传统数学启蒙"

    # Table slide
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Add title
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_shape.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = "九九乘法表"
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = PptxRGBColor(0, 51, 102)

    # Create table on slide
    rows = 9
    cols = 9
    left = Inches(0.5)
    top = Inches(1.2)
    width = Inches(9)
    height = Inches(5)

    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # Set column widths
    for i in range(cols):
        table.columns[i].width = Inches(1)

    # Fill table
    mult_data = get_multiplication_table()
    colors = [
        PptxRGBColor(200, 230, 200),  # Light green
        PptxRGBColor(200, 220, 255),  # Light blue
        PptxRGBColor(255, 220, 200),  # Light orange
        PptxRGBColor(230, 230, 230),  # Light gray
    ]

    for i in range(rows):
        for j in range(cols):
            cell = table.cell(i, j)
            if j <= i:
                cell.text = mult_data[i][j]
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.alignment = PP_ALIGN.CENTER
                    for run in paragraph.runs:
                        run.font.size = Pt(11)
                        run.font.bold = True
                # Set cell background color based on row
                fill = cell.fill
                fill.solid()
                fill.fore_color.rgb = colors[i % len(colors)]
            else:
                cell.text = ""
                cell.fill.solid()
                cell.fill.fore_color.rgb = PptxRGBColor(255, 255, 255)

    # Add summary slide
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "乘法口诀要点"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "学习要点："

    points = [
        "共 81 句口诀，从 1×1=1 到 9×9=81",
        "每行口诀数量递增：第 1 行 1 句，第 9 行 9 句",
        "乘法交换律：a×b = b×a",
        "对角线上的数是平方数：1,4,9,16,25,36,49,64,81",
        "5 的倍数结尾只有 0 或 5",
        "9 的倍数各位数字之和为 9"
    ]

    for point in points:
        p = tf.add_paragraph()
        p.text = "• " + point
        p.font.size = Pt(16)
        p.space_after = Pt(10)

    # Save presentation
    output_path = '/Users/davidli/lobsterai/project/乘法口诀表.pptx'
    prs.save(output_path)
    print(f"✓ PowerPoint saved to: {output_path}")
    return output_path

# Main execution
if __name__ == '__main__':
    print("Creating multiplication table documents...")
    print()

    # Create Word
    word_path = create_word_doc()

    # Create PowerPoint
    ppt_path = create_ppt()

    print()
    print("✅ All documents created successfully!")
